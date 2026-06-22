"""Slide layout and style helpers for PPT generation."""

import math

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Color palette
# ---------------------------------------------------------------------------
COLOR_PRIMARY    = RGBColor(0x1B, 0x36, 0x5D)   # SRM Navy
COLOR_ACCENT     = RGBColor(0xC8, 0x10, 0x2E)   # SRM Maroon
COLOR_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
COLOR_DARK_TEXT  = RGBColor(0x33, 0x33, 0x33)
COLOR_SUBTITLE   = RGBColor(0x66, 0x66, 0x66)
COLOR_CODE_BG    = RGBColor(0x1E, 0x1E, 0x1E)
COLOR_CODE_TEXT  = RGBColor(0xD4, 0xD4, 0xD4)
COLOR_HIGHLIGHT  = RGBColor(0xFF, 0xF3, 0xCD)   # soft yellow for analogy / quiz

# ---------------------------------------------------------------------------
# Dimensions
# ---------------------------------------------------------------------------
SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
FONT_FAMILY  = "Calibri"
FONT_MONO    = "Courier New"

# MSO shape-type integers (matches Office VBA constants)
_RECT         = 1
_ROUNDED_RECT = 5
_OVAL         = 9



# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def create_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs


def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _accent_bar(slide, top=False):
    y = Inches(0) if top else SLIDE_HEIGHT - Inches(0.15)
    s = slide.shapes.add_shape(_RECT, Inches(0), y, SLIDE_WIDTH, Inches(0.15))
    s.fill.solid()
    s.fill.fore_color.rgb = COLOR_ACCENT
    s.line.fill.background()


def _footer(slide, module_num, session_num, slo_num):
    tb = slide.shapes.add_textbox(
        Inches(0.5), SLIDE_HEIGHT - Inches(0.55), Inches(7), Inches(0.35))
    p = tb.text_frame.paragraphs[0]
    p.text = f"Module {module_num}  |  Session {session_num}  |  SLO {slo_num}"
    p.font.size = Pt(9)
    p.font.color.rgb = COLOR_SUBTITLE
    p.font.name = FONT_FAMILY


def _title_bar(slide, text):
    """Navy bar across the top with white title text."""
    bar = slide.shapes.add_shape(_RECT, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_PRIMARY
    bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.15), Inches(12.1), Inches(0.75))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.font.name = FONT_FAMILY


def _textbox(slide, left, top, width, height,
             text="", font_size=16, color=None, bold=False,
             italic=False, align=PP_ALIGN.LEFT, wrap=True, mono=False):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color or COLOR_DARK_TEXT
    p.font.name = FONT_MONO if mono else FONT_FAMILY
    p.alignment = align
    return tb, tf


def _filled_box(slide, left, top, width, height, fill_color, shape_type=_RECT):
    s = slide.shapes.add_shape(
        shape_type,
        Inches(left), Inches(top), Inches(width), Inches(height))
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    s.line.fill.background()
    return s


def _bullet_list(slide, left, top, width, height, items,
                 font_size=15, color=None, prefix="•"):
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{prefix} {item}" if prefix else item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color or COLOR_DARK_TEXT
        p.font.name = FONT_FAMILY
        p.space_after = Pt(8)
    return tb


def _set_notes(slide, text):
    if text:
        slide.notes_slide.notes_text_frame.text = text


def _parse_color(hex_str, fallback):
    try:
        h = hex_str.lstrip("#")
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    except Exception:
        return fallback


# ---------------------------------------------------------------------------
# 1. TITLE SLIDE
# ---------------------------------------------------------------------------

def add_title_slide(prs, slide_data, module_num, session_num, slo_num):
    slide = _blank_slide(prs)

    # Full-width navy background panel (top 60%)
    bg = slide.shapes.add_shape(_RECT, Inches(0), Inches(0), SLIDE_WIDTH, Inches(4.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_PRIMARY
    bg.line.fill.background()

    # Title
    _textbox(slide, 0.9, 0.7, 11.5, 1.8,
             text=slide_data.get("title", ""),
             font_size=36, bold=True, color=COLOR_WHITE,
             align=PP_ALIGN.LEFT)

    # SLO subtitle
    _textbox(slide, 0.9, 2.6, 11.5, 1.1,
             text=slide_data.get("subtitle", ""),
             font_size=18, color=RGBColor(0xCC, 0xCC, 0xFF),
             align=PP_ALIGN.LEFT)

    # Hook box (light area)
    hook_bg = slide.shapes.add_shape(_ROUNDED_RECT, Inches(0.8), Inches(4.8), Inches(11.7), Inches(0.9))
    hook_bg.fill.solid()
    hook_bg.fill.fore_color.rgb = COLOR_HIGHLIGHT
    hook_bg.line.color.rgb = COLOR_ACCENT
    hook_bg.line.width = Pt(1.5)
    tf = hook_bg.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = slide_data.get("hook", "")
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = COLOR_DARK_TEXT
    p.font.name = FONT_FAMILY
    p.alignment = PP_ALIGN.CENTER
    try:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except Exception:
        pass

    # Course info
    _textbox(slide, 0.9, 5.85, 11.5, 0.4,
             text=f"21CSE597T — Containers and Cloud DevOps  |  Module {module_num}  |  Session {session_num}  |  SLO {slo_num}",
             font_size=11, color=COLOR_SUBTITLE)

    _accent_bar(slide, top=False)
    _set_notes(slide, slide_data.get("speaker_notes", ""))
    return slide


# ---------------------------------------------------------------------------
# 2. WHAT SLIDE  (two_col)
# ---------------------------------------------------------------------------

def add_what_slide(prs, slide_data, module_num, session_num, slo_num):
    slide = _blank_slide(prs)
    _title_bar(slide, slide_data.get("title", "What is it?"))

    # Left column — Definition
    _textbox(slide, 0.5, 1.2, 0.9, 0.35,
             text="DEFINITION", font_size=10, bold=True, color=COLOR_ACCENT)

    _textbox(slide, 0.5, 1.55, 6.2, 2.3,
             text=slide_data.get("definition", ""),
             font_size=14, color=COLOR_DARK_TEXT)

    # Analogy box
    analogy_bg = slide.shapes.add_shape(_ROUNDED_RECT, Inches(0.5), Inches(4.0), Inches(6.2), Inches(1.4))
    analogy_bg.fill.solid()
    analogy_bg.fill.fore_color.rgb = COLOR_HIGHLIGHT
    analogy_bg.line.color.rgb = COLOR_ACCENT
    analogy_bg.line.width = Pt(1)
    tf = analogy_bg.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"💡  {slide_data.get('analogy', '')}"
    p.font.size = Pt(13)
    p.font.italic = True
    p.font.color.rgb = COLOR_DARK_TEXT
    p.font.name = FONT_FAMILY
    try:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except Exception:
        pass

    # Divider line
    div = slide.shapes.add_shape(_RECT, Inches(7.0), Inches(1.1), Inches(0.04), Inches(5.1))
    div.fill.solid()
    div.fill.fore_color.rgb = COLOR_LIGHT_GRAY
    div.line.fill.background()

    # Right column — Key Characteristics
    _textbox(slide, 7.2, 1.2, 5.7, 0.35,
             text="KEY CHARACTERISTICS", font_size=10, bold=True, color=COLOR_ACCENT)

    chars = slide_data.get("key_characteristics", [])
    _bullet_list(slide, 7.2, 1.65, 5.7, 4.8, chars, font_size=15)

    _footer(slide, module_num, session_num, slo_num)
    _accent_bar(slide, top=False)
    _set_notes(slide, slide_data.get("speaker_notes", ""))
    return slide


# ---------------------------------------------------------------------------
# 3. WHY SLIDE  (content)
# ---------------------------------------------------------------------------

def add_why_slide(prs, slide_data, module_num, session_num, slo_num):
    slide = _blank_slide(prs)
    _title_bar(slide, slide_data.get("title", "Why does it matter?"))

    # Problem statement banner
    ps_bg = slide.shapes.add_shape(_RECT, Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.8))
    ps_bg.fill.solid()
    ps_bg.fill.fore_color.rgb = RGBColor(0xE8, 0xEC, 0xF5)
    ps_bg.line.color.rgb = COLOR_PRIMARY
    ps_bg.line.width = Pt(1)
    tf = ps_bg.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"Problem: {slide_data.get('problem_statement', '')}"
    p.font.size = Pt(14)
    p.font.italic = True
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.font.name = FONT_FAMILY
    try:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except Exception:
        pass

    # Bullet points
    bullets = slide_data.get("bullet_points", [])
    _bullet_list(slide, 0.7, 2.2, 12.0, 3.5, bullets, font_size=16)

    # Industry context box
    ic_bg = slide.shapes.add_shape(_ROUNDED_RECT, Inches(0.5), Inches(5.85), Inches(12.3), Inches(0.85))
    ic_bg.fill.solid()
    ic_bg.fill.fore_color.rgb = COLOR_PRIMARY
    ic_bg.line.fill.background()
    tf = ic_bg.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"Industry: {slide_data.get('industry_context', '')}"
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_WHITE
    p.font.name = FONT_FAMILY
    try:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except Exception:
        pass

    _footer(slide, module_num, session_num, slo_num)
    _accent_bar(slide, top=False)
    _set_notes(slide, slide_data.get("speaker_notes", ""))
    return slide


# ---------------------------------------------------------------------------
# 4. HOW SLIDE  (two_col)
# ---------------------------------------------------------------------------

def add_how_slide(prs, slide_data, module_num, session_num, slo_num):
    slide = _blank_slide(prs)
    _title_bar(slide, slide_data.get("title", "How does it work?"))

    # Left column — Steps
    _textbox(slide, 0.5, 1.2, 0.6, 0.35,
             text="STEPS", font_size=10, bold=True, color=COLOR_ACCENT)

    steps = slide_data.get("steps", [])
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(6.5), Inches(5.1))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, step in enumerate(steps):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.font.name = FONT_FAMILY
        p.font.size = Pt(15)
        p.font.color.rgb = COLOR_DARK_TEXT
        p.space_after = Pt(10)
        run = p.add_run()
        run.text = f"{i + 1}. "
        run.font.bold = True
        run.font.color.rgb = COLOR_PRIMARY
        run.font.name = FONT_FAMILY
        run2 = p.add_run()
        run2.text = step
        run2.font.color.rgb = COLOR_DARK_TEXT
        run2.font.name = FONT_FAMILY

    # Divider
    div = slide.shapes.add_shape(_RECT, Inches(7.3), Inches(1.1), Inches(0.04), Inches(5.7))
    div.fill.solid()
    div.fill.fore_color.rgb = COLOR_LIGHT_GRAY
    div.line.fill.background()

    # Right column — Tools
    _textbox(slide, 7.5, 1.2, 0.6, 0.35,
             text="TOOLS", font_size=10, bold=True, color=COLOR_ACCENT)

    tools = slide_data.get("tools", [])
    ty = 1.65
    for tool in tools:
        name = tool.get("name", "")
        role = tool.get("role", "")
        tool_bg = slide.shapes.add_shape(_ROUNDED_RECT, Inches(7.5), Inches(ty), Inches(5.4), Inches(0.72))
        tool_bg.fill.solid()
        tool_bg.fill.fore_color.rgb = RGBColor(0xE8, 0xEC, 0xF5)
        tool_bg.line.color.rgb = COLOR_PRIMARY
        tool_bg.line.width = Pt(0.75)
        tf = tool_bg.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run1 = p.add_run()
        run1.text = f"{name}  "
        run1.font.bold = True
        run1.font.size = Pt(13)
        run1.font.color.rgb = COLOR_PRIMARY
        run1.font.name = FONT_FAMILY
        run2 = p.add_run()
        run2.text = role
        run2.font.size = Pt(12)
        run2.font.color.rgb = COLOR_DARK_TEXT
        run2.font.name = FONT_FAMILY
        try:
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        except Exception:
            pass
        ty += 0.85

    _footer(slide, module_num, session_num, slo_num)
    _accent_bar(slide, top=False)
    _set_notes(slide, slide_data.get("speaker_notes", ""))
    return slide


# ---------------------------------------------------------------------------
# 5. DIAGRAM SLIDE
# ---------------------------------------------------------------------------

def _draw_diagram(slide, diagram_data, ax, ay, aw, ah):
    """Render nodes and edges in the given area (all values in Inches float)."""
    nodes = diagram_data.get("nodes", [])
    edges = diagram_data.get("edges", [])
    layout_dir = diagram_data.get("layout_direction", "left_to_right")

    if not nodes:
        return

    n = len(nodes)
    NODE_W = Inches(1.9)
    NODE_H = Inches(0.6)

    # --- Position nodes ---
    positions = {}  # id -> (center_x_emu, center_y_emu)

    if layout_dir == "top_to_bottom":
        cols = min(3, n)
        rows = math.ceil(n / cols)
        col_w = Inches(aw) / cols
        row_h = Inches(ah) / rows
        for i, node in enumerate(nodes):
            cx = Inches(ax) + (i % cols) * col_w + col_w // 2
            cy = Inches(ay) + (i // cols) * row_h + row_h // 2
            positions[node["id"]] = (cx, cy)
    else:  # left_to_right
        if n <= 5:
            col_w = Inches(aw) / n
            for i, node in enumerate(nodes):
                cx = Inches(ax) + i * col_w + col_w // 2
                cy = Inches(ay) + Inches(ah) // 2
                positions[node["id"]] = (cx, cy)
        else:
            top_n = math.ceil(n / 2)
            for i, node in enumerate(nodes):
                row = 0 if i < top_n else 1
                col = i if i < top_n else i - top_n
                row_nodes = top_n if row == 0 else n - top_n
                col_w = Inches(aw) / row_nodes
                cx = Inches(ax) + col * col_w + col_w // 2
                cy = Inches(ay) + Inches(ah) * (0.3 if row == 0 else 0.75)
                positions[node["id"]] = (cx, cy)

    # --- Draw edges (thin rectangles — avoids connector XML issues in PowerPoint) ---
    for edge in edges:
        fid, tid = edge.get("from"), edge.get("to")
        if fid not in positions or tid not in positions:
            continue
        fx, fy = positions[fid]
        tx, ty = positions[tid]

        if abs(tx - fx) >= abs(ty - fy):
            # Horizontal-dominant: draw horizontal bar + small arrowhead box
            lx = min(fx, tx)
            lw = abs(tx - fx)
            ly = (fy + ty) // 2 - Inches(0.025)
            lh = Inches(0.05)
        else:
            # Vertical-dominant: draw vertical bar
            lx = (fx + tx) // 2 - Inches(0.025)
            lw = Inches(0.05)
            ly = min(fy, ty)
            lh = abs(ty - fy)

        if lw > 0 and lh > 0:
            ln = slide.shapes.add_shape(_RECT, lx, ly, lw, lh)
            ln.fill.solid()
            ln.fill.fore_color.rgb = COLOR_PRIMARY
            ln.line.fill.background()

        # Arrowhead: small right-pointing triangle using a narrow rectangle at the destination end
        if edge.get("arrow", True) and abs(tx - fx) >= abs(ty - fy):
            arrow_x = max(fx, tx) - Inches(0.15)
            arrow_y = (fy + ty) // 2 - Inches(0.08)
            arr = slide.shapes.add_shape(5, arrow_x, arrow_y, Inches(0.15), Inches(0.16))
            arr.fill.solid()
            arr.fill.fore_color.rgb = COLOR_PRIMARY
            arr.line.fill.background()

        label = edge.get("label", "")
        if label:
            lx = (fx + tx) // 2 - Inches(0.6)
            ly = (fy + ty) // 2 - Inches(0.18)
            lb = slide.shapes.add_textbox(lx, ly, Inches(1.2), Inches(0.35))
            p = lb.text_frame.paragraphs[0]
            p.text = label
            p.font.size = Pt(9)
            p.font.color.rgb = COLOR_DARK_TEXT
            p.font.name = FONT_FAMILY
            p.alignment = PP_ALIGN.CENTER

    # --- Draw nodes (on top of edges) ---
    for node in nodes:
        nid = node["id"]
        if nid not in positions:
            continue
        cx, cy = positions[nid]
        node_left = cx - NODE_W // 2
        node_top  = cy - NODE_H // 2

        node_color = _parse_color(node.get("color", "#1B365D"), COLOR_PRIMARY)
        font_color = _parse_color(node.get("font_color", "#FFFFFF"), COLOR_WHITE)
        shape_type = _OVAL if node.get("shape") in ("circle", "ellipse") else (
                     _ROUNDED_RECT if node.get("shape") == "rounded" else _RECT)

        rect = slide.shapes.add_shape(shape_type, node_left, node_top, NODE_W, NODE_H)
        rect.fill.solid()
        rect.fill.fore_color.rgb = node_color
        rect.line.color.rgb = COLOR_WHITE
        rect.line.width = Pt(1)

        tf = rect.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = node["label"]
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = font_color
        p.font.name = FONT_FAMILY
        p.alignment = PP_ALIGN.CENTER
        try:
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        except Exception:
            pass


def add_diagram_slide(prs, slide_data, module_num, session_num, slo_num):
    slide = _blank_slide(prs)
    _title_bar(slide, slide_data.get("title", "Architecture Diagram"))

    diagram_data = slide_data.get("diagram", {})
    _draw_diagram(slide, diagram_data, ax=0.5, ay=1.2, aw=12.3, ah=4.8)

    # Caption
    caption = diagram_data.get("caption", "")
    if caption:
        cap_bg = slide.shapes.add_shape(_RECT, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.45))
        cap_bg.fill.solid()
        cap_bg.fill.fore_color.rgb = COLOR_LIGHT_GRAY
        cap_bg.line.fill.background()
        tf = cap_bg.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = caption
        p.font.size = Pt(11)
        p.font.italic = True
        p.font.color.rgb = COLOR_SUBTITLE
        p.font.name = FONT_FAMILY
        p.alignment = PP_ALIGN.CENTER
        try:
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        except Exception:
            pass

    _footer(slide, module_num, session_num, slo_num)
    _accent_bar(slide, top=False)
    _set_notes(slide, slide_data.get("speaker_notes", ""))
    return slide


# ---------------------------------------------------------------------------
# 6. CODE SLIDE
# ---------------------------------------------------------------------------

def add_code_slide(prs, slide_data, module_num, session_num, slo_num):
    slide = _blank_slide(prs)
    _title_bar(slide, slide_data.get("title", "Code Example"))

    # Language badge
    lang = slide_data.get("language", "")
    if lang:
        badge = slide.shapes.add_shape(_ROUNDED_RECT, Inches(0.5), Inches(1.15), Inches(1.4), Inches(0.35))
        badge.fill.solid()
        badge.fill.fore_color.rgb = COLOR_ACCENT
        badge.line.fill.background()
        tf = badge.text_frame
        p = tf.paragraphs[0]
        p.text = lang.upper()
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.font.name = FONT_FAMILY
        p.alignment = PP_ALIGN.CENTER
        try:
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        except Exception:
            pass

    # Dark code block background
    code_bg = slide.shapes.add_shape(_RECT, Inches(0.5), Inches(1.6), Inches(8.4), Inches(4.4))
    code_bg.fill.solid()
    code_bg.fill.fore_color.rgb = COLOR_CODE_BG
    code_bg.line.fill.background()

    # Code text (monospace)
    code = slide_data.get("code_snippet", "")
    tb = slide.shapes.add_textbox(Inches(0.75), Inches(1.75), Inches(7.9), Inches(4.1))
    tf = tb.text_frame
    tf.word_wrap = False
    for i, line in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_CODE_TEXT
        p.font.name = FONT_MONO
        p.space_after = Pt(1)

    # Explanation annotations (right panel)
    _textbox(slide, 9.2, 1.15, 3.8, 0.35,
             text="ANNOTATIONS", font_size=10, bold=True, color=COLOR_ACCENT)

    explanations = slide_data.get("explanation", [])
    ey = 1.6
    for exp in explanations[:5]:
        line_ref = exp.get("line_ref", "")
        note     = exp.get("note", "")
        ann_bg = slide.shapes.add_shape(_ROUNDED_RECT, Inches(9.2), Inches(ey), Inches(3.7), Inches(0.78))
        ann_bg.fill.solid()
        ann_bg.fill.fore_color.rgb = RGBColor(0xE8, 0xEC, 0xF5)
        ann_bg.line.color.rgb = COLOR_PRIMARY
        ann_bg.line.width = Pt(0.75)
        tf = ann_bg.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run1 = p.add_run()
        run1.text = f"{line_ref}  "
        run1.font.bold = True
        run1.font.size = Pt(11)
        run1.font.color.rgb = COLOR_ACCENT
        run1.font.name = FONT_FAMILY
        run2 = p.add_run()
        run2.text = note
        run2.font.size = Pt(11)
        run2.font.color.rgb = COLOR_DARK_TEXT
        run2.font.name = FONT_FAMILY
        try:
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        except Exception:
            pass
        ey += 0.9

    _footer(slide, module_num, session_num, slo_num)
    _accent_bar(slide, top=False)
    _set_notes(slide, slide_data.get("speaker_notes", ""))
    return slide


# ---------------------------------------------------------------------------
# 7. USE CASE SLIDE  (two_col)
# ---------------------------------------------------------------------------

def add_use_case_slide(prs, slide_data, module_num, session_num, slo_num):
    slide = _blank_slide(prs)
    _title_bar(slide, slide_data.get("title", "Real-World Use Case"))

    # Company badge
    company = slide_data.get("company", "")
    if company:
        cb = slide.shapes.add_shape(_ROUNDED_RECT, Inches(0.5), Inches(1.2), Inches(3.5), Inches(0.45))
        cb.fill.solid()
        cb.fill.fore_color.rgb = COLOR_PRIMARY
        cb.line.fill.background()
        tf = cb.text_frame
        p = tf.paragraphs[0]
        p.text = company
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.font.name = FONT_FAMILY
        p.alignment = PP_ALIGN.CENTER
        try:
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        except Exception:
            pass

    # Left column labels + content
    for label, key, top in [
        ("SCENARIO", "scenario", 1.85),
        ("SOLUTION", "solution", 3.5),
    ]:
        _textbox(slide, 0.5, top, 1.1, 0.3, text=label,
                 font_size=9, bold=True, color=COLOR_ACCENT)
        _textbox(slide, 0.5, top + 0.35, 6.4, 1.15,
                 text=slide_data.get(key, ""), font_size=13)

    # Divider
    div = slide.shapes.add_shape(_RECT, Inches(7.2), Inches(1.1), Inches(0.04), Inches(5.7))
    div.fill.solid()
    div.fill.fore_color.rgb = COLOR_LIGHT_GRAY
    div.line.fill.background()

    # Outcome highlight box
    _textbox(slide, 7.4, 1.2, 0.9, 0.3, text="OUTCOME",
             font_size=9, bold=True, color=COLOR_ACCENT)
    out_bg = slide.shapes.add_shape(_ROUNDED_RECT, Inches(7.4), Inches(1.6), Inches(5.5), Inches(2.0))
    out_bg.fill.solid()
    out_bg.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
    out_bg.line.color.rgb = RGBColor(0x2E, 0x7D, 0x32)
    out_bg.line.width = Pt(1)
    tf = out_bg.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = slide_data.get("outcome", "")
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1B, 0x5E, 0x20)
    p.font.name = FONT_FAMILY
    try:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except Exception:
        pass

    # Lesson box
    _textbox(slide, 7.4, 3.75, 0.9, 0.3, text="LESSON",
             font_size=9, bold=True, color=COLOR_ACCENT)
    lesson_bg = slide.shapes.add_shape(_ROUNDED_RECT, Inches(7.4), Inches(4.1), Inches(5.5), Inches(1.5))
    lesson_bg.fill.solid()
    lesson_bg.fill.fore_color.rgb = COLOR_HIGHLIGHT
    lesson_bg.line.color.rgb = COLOR_ACCENT
    lesson_bg.line.width = Pt(1)
    tf = lesson_bg.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = slide_data.get("lesson", "")
    p.font.size = Pt(13)
    p.font.italic = True
    p.font.color.rgb = COLOR_DARK_TEXT
    p.font.name = FONT_FAMILY
    try:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except Exception:
        pass

    _footer(slide, module_num, session_num, slo_num)
    _accent_bar(slide, top=False)
    _set_notes(slide, slide_data.get("speaker_notes", ""))
    return slide


# ---------------------------------------------------------------------------
# 8. QUIZ SLIDE
# ---------------------------------------------------------------------------

def add_quiz_slide(prs, slide_data, module_num, session_num, slo_num):
    slide = _blank_slide(prs)

    # Header bar in accent color
    bar = slide.shapes.add_shape(_RECT, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_ACCENT
    bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.15), Inches(12.1), Inches(0.75))
    p = tb.text_frame.paragraphs[0]
    p.text = "Knowledge Check"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.font.name = FONT_FAMILY

    # Question box
    q_bg = slide.shapes.add_shape(_ROUNDED_RECT, Inches(0.5), Inches(1.2), Inches(12.3), Inches(1.0))
    q_bg.fill.solid()
    q_bg.fill.fore_color.rgb = COLOR_HIGHLIGHT
    q_bg.line.color.rgb = COLOR_ACCENT
    q_bg.line.width = Pt(1.5)
    tf = q_bg.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = slide_data.get("question", "")
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_DARK_TEXT
    p.font.name = FONT_FAMILY
    try:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except Exception:
        pass

    # Options (2×2 grid)
    options = slide_data.get("options", {})
    correct = slide_data.get("correct_answer", "")
    positions = [
        ("A", 0.5,  2.4),
        ("B", 6.9,  2.4),
        ("C", 0.5,  3.5),
        ("D", 6.9,  3.5),
    ]
    for letter, ox, oy in positions:
        is_correct = (letter == correct)
        opt_text = f"{letter}.  {options.get(letter, '')}"
        opt_bg = slide.shapes.add_shape(
            _ROUNDED_RECT, Inches(ox), Inches(oy), Inches(6.0), Inches(0.82))
        opt_bg.fill.solid()
        opt_bg.fill.fore_color.rgb = (
            RGBColor(0xE8, 0xF5, 0xE9) if is_correct else RGBColor(0xF5, 0xF5, 0xF5))
        opt_bg.line.color.rgb = (
            RGBColor(0x2E, 0x7D, 0x32) if is_correct else COLOR_LIGHT_GRAY)
        opt_bg.line.width = Pt(1.5 if is_correct else 0.75)
        tf = opt_bg.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = opt_text
        p.font.size = Pt(13)
        p.font.bold = is_correct
        p.font.color.rgb = (RGBColor(0x1B, 0x5E, 0x20) if is_correct else COLOR_DARK_TEXT)
        p.font.name = FONT_FAMILY
        try:
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        except Exception:
            pass

    # Challenge task box
    challenge = slide_data.get("challenge_task", "")
    if challenge:
        ch_bg = slide.shapes.add_shape(_ROUNDED_RECT, Inches(0.5), Inches(4.55), Inches(12.3), Inches(1.0))
        ch_bg.fill.solid()
        ch_bg.fill.fore_color.rgb = RGBColor(0xE8, 0xEC, 0xF5)
        ch_bg.line.color.rgb = COLOR_PRIMARY
        ch_bg.line.width = Pt(1)
        tf = ch_bg.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run1 = p.add_run()
        run1.text = "Challenge:  "
        run1.font.bold = True
        run1.font.size = Pt(13)
        run1.font.color.rgb = COLOR_PRIMARY
        run1.font.name = FONT_FAMILY
        run2 = p.add_run()
        run2.text = challenge
        run2.font.size = Pt(13)
        run2.font.color.rgb = COLOR_DARK_TEXT
        run2.font.name = FONT_FAMILY
        try:
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        except Exception:
            pass

    _footer(slide, module_num, session_num, slo_num)
    _accent_bar(slide, top=False)
    _set_notes(slide, slide_data.get("speaker_notes", ""))
    return slide


# ---------------------------------------------------------------------------
# 9. SUMMARY SLIDE
# ---------------------------------------------------------------------------

def add_summary_slide(prs, slide_data, module_num, session_num, slo_num):
    slide = _blank_slide(prs)

    # Title bar in accent
    bar = slide.shapes.add_shape(_RECT, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_ACCENT
    bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.15), Inches(12.1), Inches(0.75))
    p = tb.text_frame.paragraphs[0]
    p.text = slide_data.get("title", "Key Takeaways")
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.font.name = FONT_FAMILY

    # Numbered takeaways
    bullets = slide_data.get("bullet_points", [])
    tb2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(12.0), Inches(3.8))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    labels = ["WHAT", "WHY", "HOW", "USE CASE", "5", "6"]
    for i, point in enumerate(bullets):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        label = labels[i] if i < len(labels) else str(i + 1)
        run1 = p.add_run()
        run1.text = f"[{label}]  "
        run1.font.bold = True
        run1.font.size = Pt(16)
        run1.font.color.rgb = COLOR_PRIMARY
        run1.font.name = FONT_FAMILY
        run2 = p.add_run()
        run2.text = point
        run2.font.size = Pt(16)
        run2.font.color.rgb = COLOR_DARK_TEXT
        run2.font.name = FONT_FAMILY
        p.space_after = Pt(12)

    # CTA box
    cta = slide_data.get("call_to_action", "")
    if cta:
        cta_bg = slide.shapes.add_shape(_ROUNDED_RECT, Inches(0.6), Inches(5.4), Inches(7.5), Inches(0.75))
        cta_bg.fill.solid()
        cta_bg.fill.fore_color.rgb = COLOR_PRIMARY
        cta_bg.line.fill.background()
        tf = cta_bg.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run1 = p.add_run()
        run1.text = "Next Step:  "
        run1.font.bold = True
        run1.font.size = Pt(13)
        run1.font.color.rgb = COLOR_ACCENT
        run1.font.name = FONT_FAMILY
        run2 = p.add_run()
        run2.text = cta
        run2.font.size = Pt(13)
        run2.font.color.rgb = COLOR_WHITE
        run2.font.name = FONT_FAMILY
        try:
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        except Exception:
            pass

    # Preview teaser
    preview = slide_data.get("preview", "")
    if preview:
        _textbox(slide, 0.6, 6.3, 12.0, 0.45,
                 text=f"Up next: {preview}",
                 font_size=12, italic=True, color=COLOR_SUBTITLE)

    _footer(slide, module_num, session_num, slo_num)
    _accent_bar(slide, top=False)
    _set_notes(slide, slide_data.get("speaker_notes", ""))
    return slide


# ---------------------------------------------------------------------------
# Fallback generic content slide (for unknown slide types)
# ---------------------------------------------------------------------------

def add_content_slide(prs, slide_data, module_num, session_num, slo_num):
    slide = _blank_slide(prs)
    title   = slide_data.get("title", "")
    bullets = slide_data.get("bullet_points", [])
    notes   = slide_data.get("speaker_notes", "")

    _title_bar(slide, title)
    _bullet_list(slide, 0.8, 1.4, 12.0, 5.3, bullets, font_size=16)

    _footer(slide, module_num, session_num, slo_num)
    _accent_bar(slide, top=False)
    _set_notes(slide, notes)
    return slide
