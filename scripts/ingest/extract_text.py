"""Extract text from various document formats (PDF, DOCX, PPTX, web URLs)."""

import json
import sys
from pathlib import Path

import pdfplumber
from bs4 import BeautifulSoup
from docx import Document
from pptx import Presentation

import requests


def extract_from_pdf(file_path: str) -> str:
    """Extract text from a PDF file."""
    text_parts = []
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text_parts.append(page_text)
    return "\n\n".join(text_parts)


def extract_from_docx(file_path: str) -> str:
    """Extract text from a DOCX file."""
    doc = Document(file_path)
    text_parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            text_parts.append(para.text.strip())
    # Also extract from tables
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                text_parts.append(row_text)
    return "\n\n".join(text_parts)


def extract_from_pptx(file_path: str) -> str:
    """Extract text from a PPTX file."""
    prs = Presentation(file_path)
    text_parts = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_texts.append(text)
        if slide_texts:
            text_parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_texts))
    return "\n\n".join(text_parts)


def extract_from_url(url: str) -> str:
    """Extract text content from a web URL."""
    try:
        response = requests.get(url, timeout=30, headers={"User-Agent": "Mozilla/5.0"})
        response.raise_for_status()
        soup = BeautifulSoup(response.text, "html.parser")
        # Remove script and style elements
        for element in soup(["script", "style", "nav", "footer", "header"]):
            element.decompose()
        text = soup.get_text(separator="\n", strip=True)
        # Clean up excessive whitespace
        lines = [line.strip() for line in text.splitlines() if line.strip()]
        return "\n".join(lines)
    except Exception as e:
        return f"[Error extracting from {url}: {e}]"


def extract_text(file_path: str) -> str:
    """Extract text from a file based on its extension, or from a URL."""
    if file_path.startswith(("http://", "https://")):
        return extract_from_url(file_path)

    path = Path(file_path)
    ext = path.suffix.lower()

    extractors = {
        ".pdf": extract_from_pdf,
        ".docx": extract_from_docx,
        ".pptx": extract_from_pptx,
        ".txt": lambda p: Path(p).read_text(encoding="utf-8"),
        ".md": lambda p: Path(p).read_text(encoding="utf-8"),
    }

    extractor = extractors.get(ext)
    if extractor is None:
        return f"[Unsupported file format: {ext}]"

    return extractor(file_path)


def extract_all_sources(source_dir: str) -> dict[str, list[str]]:
    """Extract text from all source documents organized by unit.

    Expects directory structure:
        source_docs/
            unit_1/
                doc1.pdf
                doc2.docx
            unit_2/
                ...

    Returns dict mapping unit keys to lists of extracted text.
    """
    source_path = Path(source_dir)
    results = {}

    for unit_dir in sorted(source_path.iterdir()):
        if not unit_dir.is_dir() or not unit_dir.name.startswith("unit_"):
            continue

        unit_key = unit_dir.name  # e.g., "unit_1"
        texts = []

        for file in sorted(unit_dir.iterdir()):
            if file.is_file() and not file.name.startswith("."):
                text = extract_text(str(file))
                if text:
                    texts.append({"source": file.name, "text": text})

        results[unit_key] = texts

    return results


if __name__ == "__main__":
    if len(sys.argv) < 2:
        source_dir = Path(__file__).parent.parent.parent / "input" / "source_docs"
    else:
        source_dir = sys.argv[1]

    results = extract_all_sources(str(source_dir))

    # Print summary
    for unit, texts in results.items():
        total_chars = sum(len(t["text"]) for t in texts)
        print(f"{unit}: {len(texts)} sources, {total_chars} chars total", file=sys.stderr)

    print(json.dumps(results, indent=2))
