"""Generate CurriculAI hackathon presentation."""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ─────────────────────────────────────────────────────────────────
DARK       = RGBColor(0x0D, 0x1B, 0x2E)   # near-black navy
NAVY       = RGBColor(0x1B, 0x36, 0x5D)   # SRM navy
MAROON     = RGBColor(0xC8, 0x10, 0x2E)   # SRM maroon
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xB0, 0xBE, 0xC5)
CARD       = RGBColor(0x14, 0x26, 0x45)   # slightly lighter than DARK
HIGHLIGHT  = RGBColor(0xFF, 0xD7, 0x00)   # gold for key numbers
TEAL       = RGBColor(0x00, 0xB4, 0xD8)   # tech accent

W = Inches(13.33)   # slide width  (16:9)
H = Inches(7.5)     # slide height
EMU = 914400        # 1 inch in EMU


# ── Helpers ──────────────────────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank(prs):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)


def rect(slide, l, t, w, h, fill, alpha=None, line=None, line_w=Pt(0)):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(1, l, t, w, h)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape


def txt(slide, text, l, t, w, h, size, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True, italic=False):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return box


def bg(slide, color=DARK):
    rect(slide, 0, 0, W, H, color)


def top_bar(slide, title_text, subtitle=None):
    """Navy bar across top with slide title."""
    bar_h = Inches(1.1)
    rect(slide, 0, 0, W, bar_h, NAVY)
    # maroon left accent
    rect(slide, 0, 0, Inches(0.08), bar_h, MAROON)
    txt(slide, title_text,
        Inches(0.25), Inches(0.15), Inches(10), Inches(0.7),
        28, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle,
            Inches(0.25), Inches(0.72), Inches(10), Inches(0.35),
            13, color=LIGHT_GRAY)


def slide_number(slide, n):
    txt(slide, str(n),
        W - Inches(0.6), H - Inches(0.45), Inches(0.5), Inches(0.4),
        11, color=LIGHT_GRAY, align=PP_ALIGN.RIGHT)


def card(slide, l, t, w, h, fill=CARD, radius=None):
    return rect(slide, l, t, w, h, fill, line=RGBColor(0x2A, 0x3F, 0x60), line_w=Pt(0.5))


def stat_block(slide, l, t, number, label, num_color=HIGHLIGHT):
    card(slide, l, t, Inches(2.5), Inches(1.55))
    txt(slide, number,
        l + Inches(0.15), t + Inches(0.1), Inches(2.2), Inches(0.85),
        40, bold=True, color=num_color, align=PP_ALIGN.CENTER)
    txt(slide, label,
        l + Inches(0.1), t + Inches(0.95), Inches(2.3), Inches(0.55),
        12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


# ── Slide 1 — Cover ──────────────────────────────────────────────────────────
def slide_cover(prs):
    s = blank(prs)
    bg(s, DARK)

    # Left vertical maroon stripe
    rect(s, 0, 0, Inches(0.12), H, MAROON)

    # Decorative geometric: big faded navy rectangle (right side)
    rect(s, W - Inches(4.5), Inches(0.5), Inches(4.2), Inches(6.5), CARD)
    rect(s, W - Inches(4.3), Inches(0.3), Inches(4.0), Inches(6.5), NAVY)

    # AI label top-right corner of the box
    txt(s, "AI", W - Inches(3.8), Inches(0.6), Inches(3.5), Inches(5.5),
        160, bold=True, color=RGBColor(0x1E, 0x35, 0x5F), align=PP_ALIGN.CENTER)

    # Main title
    txt(s, "CurriculAI",
        Inches(0.3), Inches(1.5), Inches(8), Inches(1.8),
        72, bold=True, color=WHITE)

    # Maroon underline bar under title
    rect(s, Inches(0.3), Inches(3.15), Inches(3.6), Inches(0.07), MAROON)

    # Tagline
    txt(s, "From SLO to Slide — in Seconds.",
        Inches(0.3), Inches(3.3), Inches(8.5), Inches(0.7),
        22, color=LIGHT_GRAY, italic=True)

    # Description
    txt(s, "AI-powered course content automation for faculty.\nNo more manual slides. No more wasted weekends.",
        Inches(0.3), Inches(4.0), Inches(8), Inches(1.2),
        15, color=LIGHT_GRAY)

    # Bottom strip
    rect(s, 0, H - Inches(0.8), W, Inches(0.8), NAVY)
    txt(s, "SRM Institute of Science and Technology  ·  Hackathon 2025  ·  21CSE597T — Containers & Cloud DevOps",
        Inches(0.25), H - Inches(0.68), W - Inches(0.5), Inches(0.55),
        11, color=LIGHT_GRAY)


# ── Slide 2 — The Problem ────────────────────────────────────────────────────
def slide_problem(prs):
    s = blank(prs)
    bg(s)
    top_bar(s, "The Problem", "What we kept hearing from faculty")

    # Big quote block
    card(s, Inches(0.3), Inches(1.3), Inches(12.73), Inches(1.5))
    txt(s, '"I spend an entire weekend every semester just making slides. Copy from the textbook, format it,\nadd diagrams, make it look branded… and then I do it again for every module."',
        Inches(0.55), Inches(1.4), Inches(12.2), Inches(1.3),
        14, italic=True, color=RGBColor(0xFF, 0xE0, 0x80))

    # 3 pain point cards
    pain = [
        ("⏱  40–60 hrs", "wasted per faculty\nmember per semester\ncreating slides"),
        ("📋  Repetitive Work", "SLO documents already\nexist — slides are just\na manual reformatting"),
        ("🎨  Inconsistent Output", "Every faculty has their\nown format. No standard.\nNo branding. No structure."),
    ]
    x_start = Inches(0.3)
    gap = Inches(0.25)
    cw  = Inches(4.0)
    for i, (heading, body) in enumerate(pain):
        cx = x_start + i * (cw + gap)
        card(s, cx, Inches(2.95), cw, Inches(2.85))
        # maroon top accent on each card
        rect(s, cx, Inches(2.95), cw, Inches(0.07), MAROON)
        txt(s, heading,
            cx + Inches(0.2), Inches(3.1), cw - Inches(0.3), Inches(0.55),
            15, bold=True, color=WHITE)
        txt(s, body,
            cx + Inches(0.2), Inches(3.7), cw - Inches(0.3), Inches(1.8),
            13, color=LIGHT_GRAY)

    slide_number(s, 2)


# ── Slide 3 — The Solution ───────────────────────────────────────────────────
def slide_solution(prs):
    s = blank(prs)
    bg(s)
    top_bar(s, "The Solution — CurriculAI", "One upload. A full semester of content.")

    # Left: description
    txt(s, "What if faculty only had to write what they already write?",
        Inches(0.3), Inches(1.25), Inches(5.8), Inches(0.7),
        17, bold=True, color=WHITE)
    txt(s, (
        "Faculty already create SLO (Student Learning Outcome) documents for every course. "
        "CurriculAI reads those documents and automatically generates:\n\n"
        "  •  A 9-slide branded presentation for every SLO\n"
        "  •  A consolidated PDF learning material per module\n"
        "  •  Everything structured, formatted, and ready to upload"
    ),
        Inches(0.3), Inches(2.0), Inches(5.8), Inches(3.5),
        13, color=LIGHT_GRAY)

    # Right: outcome card
    card(s, Inches(6.4), Inches(1.2), Inches(6.6), Inches(5.5))
    rect(s, Inches(6.4), Inches(1.2), Inches(6.6), Inches(0.08), MAROON)

    txt(s, "What comes out:", Inches(6.6), Inches(1.35), Inches(6.2), Inches(0.5),
        14, bold=True, color=MAROON)

    outcomes = [
        ("90", "Branded PPT Presentations"),
        ("5",  "Module-level PDF Learning Materials"),
        ("9",  "Slide types per SLO  (What / Why / How / Diagram\n  / Code / Use Case / Quiz / Summary)"),
        ("1",  "Click from faculty to trigger it all"),
    ]
    y = Inches(1.9)
    for num, label in outcomes:
        txt(s, num,  Inches(6.6), y, Inches(1.1), Inches(0.65), 32, bold=True, color=HIGHLIGHT)
        txt(s, label, Inches(7.8), y + Inches(0.1), Inches(5.0), Inches(0.65), 12, color=WHITE)
        rect(s, Inches(6.6), y + Inches(0.7), Inches(6.0), Inches(0.01), CARD)
        y += Inches(1.1)

    slide_number(s, 3)


# ── Slide 4 — How It Works ───────────────────────────────────────────────────
def slide_how(prs):
    s = blank(prs)
    bg(s)
    top_bar(s, "How It Works", "Five steps. Fully automated.")

    steps = [
        ("01", "Faculty Uploads\nSLO Document", "PDF or DOCX with\ncourse objectives via\nour React web UI"),
        ("02", "Gemini AI\nGenerates Content", "Google Gemini Flash\nreads the SLOs and\nbuilds structured JSON\nfor every session"),
        ("03", "n8n Orchestrates\nthe Pipeline", "n8n workflow triggers\nparsing, generation,\nand build steps\nautomatically"),
        ("04", "PPTs & PDFs\nare Rendered", "python-pptx builds\nbranded slides.\nReportLab builds\nthe PDF narrative"),
        ("05", "Faculty Downloads\nor Portal Gets Updated", "Files are zipped\nper module. Portal\nupload is the\nnext step."),
    ]

    sw = Inches(2.3)
    gap = Inches(0.18)
    sx = Inches(0.25)
    sy = Inches(1.25)
    sh = Inches(4.8)

    for i, (num, title, body) in enumerate(steps):
        cx = sx + i * (sw + gap)
        # card
        card(s, cx, sy, sw, sh)
        # colored top strip per step
        colors = [TEAL, MAROON, RGBColor(0xFF, 0xA5, 0x00), NAVY, RGBColor(0x2E, 0xCC, 0x71)]
        rect(s, cx, sy, sw, Inches(0.08), colors[i])
        # step number
        txt(s, num, cx + Inches(0.1), sy + Inches(0.15), sw, Inches(0.55),
            28, bold=True, color=colors[i], align=PP_ALIGN.LEFT)
        # title
        txt(s, title, cx + Inches(0.1), sy + Inches(0.65), sw - Inches(0.15), Inches(0.9),
            12, bold=True, color=WHITE)
        # body
        txt(s, body, cx + Inches(0.1), sy + Inches(1.55), sw - Inches(0.15), Inches(2.8),
            11, color=LIGHT_GRAY)

        # Arrow between cards (except last)
        if i < len(steps) - 1:
            ax = cx + sw + Inches(0.02)
            txt(s, "→", ax, sy + Inches(2.1), Inches(0.17), Inches(0.5),
                14, bold=True, color=MAROON, align=PP_ALIGN.CENTER)

    slide_number(s, 4)


# ── Slide 5 — Tech Stack ─────────────────────────────────────────────────────
def slide_tech(prs):
    s = blank(prs)
    bg(s)
    top_bar(s, "Tech Stack", "What we built it with — and why")

    techs = [
        ("Google\nGemini Flash", "The brain. Takes SLO text and returns\nstructured JSON for every slide — What,\nWhy, How, Diagram, Code, Quiz, Summary.",
         RGBColor(0x42, 0x85, 0xF4)),
        ("n8n\nWorkflow", "Visual automation that ties everything\ntogether. Triggers the Python pipeline,\nhandles retries, no glue code needed.",
         RGBColor(0xFF, 0x69, 0x00)),
        ("FastAPI +\nPython", "Pipeline backbone. Parses SLOs, calls\nGemini, checkpoints progress after every\nSLO so crashes don't lose your work.",
         RGBColor(0x00, 0x9D, 0x57)),
        ("React +\nVite", "Faculty-facing UI. Drag-drop upload,\nmodule selector, live terminal-style\nprogress log, per-module zip downloads.",
         TEAL),
        ("python-pptx\n+ ReportLab", "Renders Gemini's JSON into actual .pptx\nfiles with SRM branding and PDF learning\nmaterials per module.",
         MAROON),
        ("Checkpoint\nSystem", "After every SLO, progress is saved.\nIf the API rate-limits mid-run,\nnext run resumes exactly where it left off.",
         RGBColor(0x9B, 0x59, 0xB6)),
    ]

    cols = 3
    cw = Inches(4.0)
    ch = Inches(2.3)
    gap_x = Inches(0.27)
    gap_y = Inches(0.2)
    sx = Inches(0.3)
    sy = Inches(1.25)

    for i, (name, desc, color) in enumerate(techs):
        row = i // cols
        col = i % cols
        cx = sx + col * (cw + gap_x)
        cy = sy + row * (ch + gap_y)
        card(s, cx, cy, cw, ch)
        rect(s, cx, cy, Inches(0.08), ch, color)
        txt(s, name, cx + Inches(0.2), cy + Inches(0.15), Inches(1.8), Inches(0.75),
            13, bold=True, color=color)
        txt(s, desc, cx + Inches(0.2), cy + Inches(0.85), cw - Inches(0.35), Inches(1.35),
            11, color=LIGHT_GRAY)

    slide_number(s, 5)


# ── Slide 6 — What We Built (Demo / Results) ─────────────────────────────────
def slide_results(prs):
    s = blank(prs)
    bg(s)
    top_bar(s, "What We Actually Built", "Not a prototype. A working pipeline.")

    # 4 big stats
    stats = [
        ("90",   "PPT Presentations\nGenerated"),
        ("5",    "Complete Modules\n(Kubernetes, Docker, CI/CD…)"),
        ("18",   "Sessions per Module\nAutomated"),
        ("100%", "Free Tier API\n(20 requests/day — we made it work)"),
    ]
    sx = Inches(0.3)
    sw = Inches(2.9)
    gap = Inches(0.27)
    sy = Inches(1.3)
    colors = [HIGHLIGHT, RGBColor(0x00, 0xFF, 0x99), TEAL, MAROON]
    for i, (num, label) in enumerate(stats):
        cx = sx + i * (sw + gap)
        card(s, cx, sy, sw, Inches(1.6))
        txt(s, num, cx, sy + Inches(0.05), sw, Inches(0.9),
            44, bold=True, color=colors[i], align=PP_ALIGN.CENTER)
        txt(s, label, cx + Inches(0.1), sy + Inches(0.95), sw - Inches(0.2), Inches(0.65),
            11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    # Bottom: what's inside each PPT
    txt(s, "Every auto-generated PPT contains:",
        Inches(0.3), Inches(3.1), Inches(6), Inches(0.5),
        14, bold=True, color=WHITE)

    slides_in_ppt = [
        "Title Slide", "What Is It?", "Why Does It Matter?",
        "How It Works (Steps)", "Architecture Diagram", "Live Code Example",
        "Real-World Use Case", "Quiz + Challenge", "Key Takeaways",
    ]
    for i, name in enumerate(slides_in_ppt):
        col = i % 3
        row = i // 3
        cx = Inches(0.4) + col * Inches(4.25)
        cy = Inches(3.65) + row * Inches(0.68)
        card(s, cx, cy, Inches(4.0), Inches(0.58))
        rect(s, cx, cy, Inches(0.07), Inches(0.58), MAROON)
        txt(s, f"{'0' + str(i+1) if i < 9 else str(i+1)}  {name}",
            cx + Inches(0.18), cy + Inches(0.1), Inches(3.7), Inches(0.45),
            12, color=WHITE)

    slide_number(s, 6)


# ── Slide 7 — Impact & What's Next ───────────────────────────────────────────
def slide_impact(prs):
    s = blank(prs)
    bg(s)
    top_bar(s, "Impact & What's Next", "Where this goes from here")

    # Left: Impact
    txt(s, "Who this helps today",
        Inches(0.3), Inches(1.25), Inches(6), Inches(0.55),
        16, bold=True, color=MAROON)

    impacts = [
        ("🎓  Faculty", "Upload one document. Walk away with a\ncomplete set of slides for the whole semester."),
        ("🏫  Institutions", "Consistent, branded, structured content\nacross all courses. No chasing faculty for slides."),
        ("👩‍💻  Students", "Better-quality, more structured learning\nmaterial in their hands faster."),
    ]
    y = Inches(1.9)
    for icon_label, desc in impacts:
        card(s, Inches(0.3), y, Inches(5.9), Inches(1.0))
        txt(s, icon_label, Inches(0.5), y + Inches(0.1), Inches(2.0), Inches(0.5),
            13, bold=True, color=WHITE)
        txt(s, desc, Inches(0.5), y + Inches(0.5), Inches(5.5), Inches(0.45),
            11, color=LIGHT_GRAY)
        y += Inches(1.12)

    # Divider
    rect(s, Inches(6.5), Inches(1.2), Inches(0.04), Inches(5.5), MAROON)

    # Right: Roadmap
    txt(s, "What's coming next",
        Inches(6.7), Inches(1.25), Inches(6), Inches(0.55),
        16, bold=True, color=TEAL)

    roadmap = [
        ("Portal Upload", "Auto-upload generated content\ndirectly to the e-curricula portal", "In Progress"),
        ("Multi-Course Support", "Faculty picks their course code.\nSystem configures itself.", "Planned"),
        ("Source Doc Ingestion", "Feed actual textbook PDFs.\nGemini grounds content in real material.", "Planned"),
        ("Faculty Dashboard", "Track generation history,\nregenerate specific sessions.", "Planned"),
    ]
    y = Inches(1.9)
    for title, desc, status in roadmap:
        status_color = RGBColor(0xFF, 0xA5, 0x00) if status == "In Progress" else RGBColor(0x5B, 0x8D, 0xB8)
        card(s, Inches(6.7), y, Inches(6.3), Inches(0.95))
        rect(s, Inches(6.7), y, Inches(0.07), Inches(0.95), status_color)
        txt(s, title, Inches(6.95), y + Inches(0.05), Inches(4.0), Inches(0.4),
            13, bold=True, color=WHITE)
        txt(s, desc, Inches(6.95), y + Inches(0.45), Inches(4.2), Inches(0.45),
            11, color=LIGHT_GRAY)
        txt(s, status, Inches(11.3), y + Inches(0.2), Inches(1.5), Inches(0.4),
            10, color=status_color, align=PP_ALIGN.RIGHT)
        y += Inches(1.08)

    slide_number(s, 7)


# ── Slide 8 — Thank You / Demo ───────────────────────────────────────────────
def slide_close(prs):
    s = blank(prs)
    bg(s, DARK)

    # Left accent stripe
    rect(s, 0, 0, Inches(0.12), H, MAROON)

    # Big decorative background element
    rect(s, W - Inches(5.5), 0, Inches(5.5), H, CARD)
    txt(s, "✓", W - Inches(4.8), Inches(0.5), Inches(4.5), Inches(5.5),
        200, bold=True, color=RGBColor(0x14, 0x26, 0x45), align=PP_ALIGN.CENTER)

    # Main message
    txt(s, "It works.", Inches(0.3), Inches(1.4), Inches(7.5), Inches(1.1),
        60, bold=True, color=WHITE)
    rect(s, Inches(0.3), Inches(2.45), Inches(3.5), Inches(0.07), MAROON)

    txt(s, (
        "90 branded PPTs. 5 PDFs. All generated\n"
        "from a single SLO document using Gemini AI\n"
        "and an n8n-orchestrated Python pipeline."
    ),
        Inches(0.3), Inches(2.6), Inches(7.8), Inches(1.5),
        16, color=LIGHT_GRAY)

    txt(s, "Live at:", Inches(0.3), Inches(4.2), Inches(3), Inches(0.45),
        13, color=LIGHT_GRAY)
    txt(s, "http://localhost:5173", Inches(0.3), Inches(4.6), Inches(5.5), Inches(0.55),
        16, bold=True, color=TEAL)

    # Bottom
    rect(s, 0, H - Inches(0.8), W, Inches(0.8), NAVY)
    txt(s, "CurriculAI  ·  SRM Institute  ·  Hackathon 2025",
        Inches(0.25), H - Inches(0.65), W - Inches(0.5), Inches(0.5),
        12, color=LIGHT_GRAY)


# ── Main ─────────────────────────────────────────────────────────────────────
def main():
    prs = new_prs()
    slide_cover(prs)
    slide_problem(prs)
    slide_solution(prs)
    slide_how(prs)
    slide_tech(prs)
    slide_results(prs)
    slide_impact(prs)
    slide_close(prs)

    out = Path(__file__).parent / "CurriculAI_Hackathon_Presentation.pptx"
    prs.save(str(out))
    print(f"Saved: {out}")
    return out


if __name__ == "__main__":
    main()
